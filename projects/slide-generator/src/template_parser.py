"""
Template parser — extract slide masters, layouts, placeholders, and theme info
from a .pptx template file.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree


NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


@dataclass
class PlaceholderInfo:
    """Information about a single placeholder in a layout."""
    idx: int
    name: str
    type: str  # title, body, date, footer, slideNumber, etc.
    left: float  # inches
    top: float
    width: float
    height: float
    font_name: Optional[str] = None
    font_size: Optional[float] = None  # pt
    font_color: Optional[str] = None  # hex
    bold: Optional[bool] = None
    alignment: Optional[str] = None


@dataclass
class LayoutInfo:
    """Information about a single layout."""
    name: str
    type: str  # inferred type: title, section, content, quote, data, end, other
    placeholders: list[PlaceholderInfo] = field(default_factory=list)
    background_color: Optional[str] = None


@dataclass
class ThemeInfo:
    """Extracted theme information."""
    accent_colors: list[str] = field(default_factory=list)
    font_major: Optional[str] = None  # heading font
    font_minor: Optional[str] = None  # body font


@dataclass
class ShapeInfo:
    """Info about a shape/icon on a slide."""
    name: str
    left: float
    top: float
    width: float
    height: float
    rotation: Optional[float] = None
    shape_type: str = "unknown"  # auto_shape, picture, group, etc.


@dataclass
class TemplateDNA:
    """Complete extracted DNA of a template."""
    slide_width: float  # inches
    slide_height: float
    layouts: list[LayoutInfo] = field(default_factory=list)
    theme: ThemeInfo = field(default_factory=ThemeInfo)
    slide_count: int = 0
    template_shapes: list[ShapeInfo] = field(default_factory=list)

    def get_layout_by_name(self, name: str) -> Optional[LayoutInfo]:
        for l in self.layouts:
            if l.name == name:
                return l
        return None

    def get_layout_by_type(self, type: str) -> Optional[LayoutInfo]:
        for l in self.layouts:
            if l.type == type:
                return l
        return None

    def get_layout_by_index(self, idx: int) -> Optional[LayoutInfo]:
        if 0 <= idx < len(self.layouts):
            return self.layouts[idx]
        return None


def _emu_to_inches(emu: int) -> float:
    return float(emu) / 914400


def _get_color_from_solid_fill(fill_elem) -> Optional[str]:
    """Extract hex color from a solid fill element."""
    srgb = fill_elem.find(".//a:srgbClr", NS)
    if srgb is not None:
        return f"#{srgb.get('val')}"
    scheme = fill_elem.find(".//a:schemeClr", NS)
    if scheme is not None:
        return f"scheme:{scheme.get('val')}"
    return None


def _infer_layout_type(name: str, placeholders: list[PlaceholderInfo]) -> str:
    """Infer the semantic type of a layout from its name and placeholders."""
    name_lower = name.lower()

    # Check for known patterns in the name
    if any(kw in name_lower for kw in ["title", "cover", "1_", "section 1"]):
        return "title"
    if any(kw in name_lower for kw in ["section", "divider", "section_header"]):
        return "section"
    if any(kw in name_lower for kw in ["quote", "citation", "testimonial"]):
        return "quote"
    if any(kw in name_lower for kw in ["data", "chart", "stat", "numbers"]):
        return "data"
    if any(kw in name_lower for kw in ["end", "closing", "thank", "last"]):
        return "end"
    if any(kw in name_lower for kw in ["icon", "legend", "gallery"]):
        return "icons"

    # Heuristic: if it has a large title placeholder + body, it's content
    has_title = any(p.type in ("title", "ctrTitle", "subTitle") for p in placeholders)
    has_body = any(p.type == "body" for p in placeholders)

    if has_title and has_body:
        return "content"
    if has_title:
        return "content"

    return "other"


def _extract_placeholders(layout) -> list[PlaceholderInfo]:
    """Extract placeholder info from a slide layout."""
    placeholders = []
    for ph in layout.placeholders:
        try:
            name = ph.name or ""
            left = _emu_to_inches(ph.left) if ph.left is not None else 0
            top = _emu_to_inches(ph.top) if ph.top is not None else 0
            width = _emu_to_inches(ph.width) if ph.width is not None else 0
            height = _emu_to_inches(ph.height) if ph.height is not None else 0

            font_name = None
            font_size = None
            font_color = None
            bold = None
            alignment = None

            # Try to get text formatting from the placeholder's XML
            txBody = ph._element.find(".//p:txBody", NS)
            if txBody is not None:
                rPr = txBody.find(".//a:rPr", NS)
                if rPr is not None:
                    typeface = rPr.get("typeface") or rPr.get("latin", {}).get("typeface")
                    if rPr.get("typeface"):
                        font_name = rPr.get("typeface")
                    sz = rPr.get("sz")
                    if sz:
                        font_size = float(sz) / 100
                    b = rPr.get("b")
                    if b:
                        bold = b == "1"

                    # Color
                    solid_fill = rPr.find(".//a:solidFill", NS)
                    if solid_fill is not None:
                        font_color = _get_color_from_solid_fill(solid_fill)

            placeholders.append(PlaceholderInfo(
                idx=ph.placeholder_format.idx if ph.placeholder_format else 0,
                name=name,
                type=ph.placeholder_format.type if ph.placeholder_format else "unknown",
                left=left,
                top=top,
                width=width,
                height=height,
                font_name=font_name,
                font_size=font_size,
                font_color=font_color,
                bold=bold,
                alignment=alignment,
            ))
        except Exception:
            continue

    return placeholders


def _extract_theme(prs: Presentation) -> ThemeInfo:
    """Extract theme info (colors, fonts) from the presentation."""
    theme = ThemeInfo()

    try:
        # Walk the XML to find the theme element
        for rel in prs.part.rels.values():
            if "theme" in rel.reltype:
                theme_elem = rel.target_part.element
                # Extract color scheme
                clrScheme = theme_elem.find(".//a:clrScheme", NS)
                if clrScheme is not None:
                    for child in clrScheme:
                        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                        srgb = child.find("a:srgbClr", NS)
                        if srgb is not None:
                            theme.accent_colors.append(f"#{srgb.get('val')}")

                # Extract major font (headings)
                major_font = theme_elem.find(".//a:majorFont/a:latin", NS)
                if major_font is not None:
                    theme.font_major = major_font.get("typeface")

                # Extract minor font (body)
                minor_font = theme_elem.find(".//a:minorFont/a:latin", NS)
                if minor_font is not None:
                    theme.font_minor = minor_font.get("typeface")

                break
    except Exception:
        pass

    return theme


def _get_background_color(layout) -> Optional[str]:
    """Extract background color from a layout."""
    try:
        bg = layout._element.find(".//p:bg", NS)
        if bg is not None:
            solid_fill = bg.find(".//a:solidFill", NS)
            if solid_fill is not None:
                return _get_color_from_solid_fill(solid_fill)
    except Exception:
        pass
    return None


def parse_template(template_path: str | Path) -> TemplateDNA:
    """
    Parse a .pptx template and return its complete DNA.

    Args:
        template_path: Path to the .pptx template file.

    Returns:
        TemplateDNA with all extracted information.
    """
    prs = Presentation(str(template_path))

    dna = TemplateDNA(
        slide_width=_emu_to_inches(prs.slide_width),
        slide_height=_emu_to_inches(prs.slide_height),
        slide_count=len(prs.slides),
    )

    # Extract theme
    dna.theme = _extract_theme(prs)

    # Extract layouts from slide layouts
    for i, layout in enumerate(prs.slide_layouts):
        try:
            name = layout.name or f"layout_{i}"
            placeholders = _extract_placeholders(layout)
            layout_type = _infer_layout_type(name, placeholders)
            bg_color = _get_background_color(layout)

            dna.layouts.append(LayoutInfo(
                name=name,
                type=layout_type,
                placeholders=placeholders,
                background_color=bg_color,
            ))
        except Exception:
            continue

    # Extract shapes/icons from template slides for reference
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                shape_type = "unknown"
                if hasattr(shape, "shape_type"):
                    shape_type = str(shape.shape_type)

                dna.template_shapes.append(ShapeInfo(
                    name=shape.name,
                    left=_emu_to_inches(shape.left) if shape.left is not None else 0,
                    top=_emu_to_inches(shape.top) if shape.top is not None else 0,
                    width=_emu_to_inches(shape.width) if shape.width is not None else 0,
                    height=_emu_to_inches(shape.height) if shape.height is not None else 0,
                    rotation=shape.rotation if hasattr(shape, "rotation") and shape.rotation else None,
                    shape_type=shape_type,
                ))
            except Exception:
                continue

    return dna


def list_layouts(template_path: str | Path) -> list[dict]:
    """Return a human-readable summary of available layouts."""
    dna = parse_template(template_path)
    result = []
    for layout in dna.layouts:
        result.append({
            "name": layout.name,
            "type": layout.type,
            "placeholders": len(layout.placeholders),
            "ph_details": [
                {"name": p.name, "type": p.type} for p in layout.placeholders
            ],
        })
    return result
