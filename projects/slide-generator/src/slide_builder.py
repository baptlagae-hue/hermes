"""
Slide builder — the core orchestrator that takes a template, plan, and content,
and builds the final presentation.
"""

from __future__ import annotations

from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu

from template_parser import parse_template, TemplateDNA
from plan_parser import Plan, SlideEntry
from content_parser import Content, SlideContent
from renderers.slides import (
    render_title_slide,
    render_section_slide,
    render_content_slide,
    render_quote_slide,
    render_data_slide,
    render_end_slide,
    render_icons_slide,
)


class SlideBuilder:
    """Builds a presentation from template + plan + content."""

    def __init__(self, template_path: str | Path):
        self.template_path = Path(template_path)
        self.dna = parse_template(self.template_path)
        self._theme_font = self.dna.theme.font_minor or self.dna.theme.font_major

    def build(self, plan: Plan, content: Optional[Content] = None,
              output_path: Optional[str | Path] = None,
              include_icons_slide: bool = True) -> Path:
        """
        Build the final presentation.

        Args:
            plan: The slide plan.
            content: Optional structured content.
            output_path: Where to save the .pptx. Auto-named if None.
            include_icons_slide: Whether to add an icons reference slide.

        Returns:
            Path to the generated .pptx file.
        """
        # Open template as base
        prs = Presentation(str(self.template_path))

        # Remove existing slides (keep only masters/layouts)
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
            prs.part.drop_rel(rId)
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

        # Build a lookup: content by slide_ref
        content_by_ref = {}
        if content:
            for sc in content.slides:
                content_by_ref[sc.slide_ref] = sc

        # Render each slide from the plan
        for i, slide_entry in enumerate(plan.slides):
            self._render_slide(prs, slide_entry, content_by_ref, i)

        # Add icons reference slide at the end
        if include_icons_slide and self.dna.template_shapes:
            icons_layout = self._find_layout_by_type("icons")
            if icons_layout:
                render_icons_slide(
                    prs, icons_layout,
                    title="Template Shapes Reference",
                    dna=self.dna,
                    theme_font=self._theme_font,
                )

        # Determine output path
        if output_path is None:
            safe_title = self._safe_filename(plan.title)
            output_path = Path(f"{safe_title}.pptx")
        else:
            output_path = Path(output_path)

        prs.save(str(output_path))
        return output_path

    def _render_slide(self, prs, slide_entry: SlideEntry,
                      content_by_ref: dict, index: int):
        """Render a single slide based on its type."""
        slide_type = slide_entry.type
        slide_content = content_by_ref.get(slide_entry.title) or content_by_ref.get(str(index))

        # Find the best layout
        layout = self._find_layout_by_type(slide_type)

        if layout is None:
            # Fallback: use the first available layout
            if self.dna.layouts:
                layout = self.dna.layouts[0]
            else:
                # Absolute fallback: use Presentation's default
                prs_obj = prs  # type: ignore
                fallback = prs_obj.slide_layouts[0] if prs_obj.slide_layouts else None
                if fallback is None:
                    raise RuntimeError("No layouts available in template!")
                layout_ref = fallback
                # We need to use the actual python-pptx layout object here
                from pptx.util import Inches as I
                slide = prs_obj.slides.add_slide(fallback)
                return

        # Get the actual pptx SlideLayout object
        prs_obj = prs
        layout_obj = self._get_pptx_layout(prs_obj, layout.name)
        if layout_obj is None:
            # Fallback to first layout
            if prs_obj.slide_layouts:
                layout_obj = prs_obj.slide_layouts[0]
            else:
                raise RuntimeError(f"Cannot find layout '{layout.name}'")

        # Extract content details
        title = slide_entry.title
        subtitle = slide_entry.subtitle
        bullets = None
        body_text = None
        quote_text = None
        quote_author = None
        data_points = []
        chart_type = "bar"

        if slide_content:
            title = slide_content.title or title
            subtitle = slide_content.subtitle or subtitle

            if slide_content.lists:
                bullets = slide_content.lists[0].items
            if slide_content.text_blocks:
                for tb in slide_content.text_blocks:
                    if tb.bullets:
                        bullets = tb.bullets
                    elif tb.text:
                        body_text = tb.text

            if slide_content.quotes:
                quote_text = slide_content.quotes[0].text
                quote_author = slide_content.quotes[0].author

            if slide_content.data:
                data_points = slide_content.data[0].data_points
                chart_type = slide_content.data[0].chart_type

        # Render by type
        if slide_type == "title":
            render_title_slide(prs_obj, layout_obj, title, subtitle, self._theme_font)
        elif slide_type == "section":
            render_section_slide(prs_obj, layout_obj, title, subtitle, self._theme_font)
        elif slide_type == "content":
            render_content_slide(prs_obj, layout_obj, title, bullets, body_text, self._theme_font)
        elif slide_type == "quote":
            txt = quote_text or "Quote text here"
            auth = quote_author
            render_quote_slide(prs_obj, layout_obj, txt, auth, self._theme_font)
        elif slide_type == "data":
            render_data_slide(prs_obj, layout_obj, title, data_points, chart_type, None, self._theme_font)
        elif slide_type == "end":
            render_end_slide(prs_obj, layout_obj, title, subtitle, self._theme_font)
        elif slide_type == "icons":
            render_icons_slide(prs_obj, layout_obj, title, self.dna, self._theme_font)
        else:
            render_content_slide(prs_obj, layout_obj, title, bullets, body_text, self._theme_font)

    def _find_layout_by_type(self, slide_type: str):
        """Find a layout matching the requested slide type."""
        if not self.dna.layouts:
            return None

        # Try exact type match
        for layout in self.dna.layouts:
            if layout.type == slide_type:
                return layout

        # Try partial name match
        name_map = {
            "title": ["title", "cover"],
            "section": ["section", "divider"],
            "content": ["content", "body", "default", "blank"],
            "quote": ["quote", "citation"],
            "data": ["data", "chart", "number"],
            "end": ["end", "closing", "thank"],
            "icons": ["icon", "legend"],
        }
        keywords = name_map.get(slide_type, [slide_type])
        for layout in self.dna.layouts:
            name_lower = layout.name.lower()
            if any(kw in name_lower for kw in keywords):
                return layout

        # Ultimate fallback: first layout
        return self.dna.layouts[0]

    def _get_pptx_layout(self, prs, layout_name: str):
        """Get a pptx SlideLayout object by its name."""
        for layout in prs.slide_layouts:
            if layout.name == layout_name:
                return layout
        return None

    def _safe_filename(self, title: str) -> str:
        """Convert a title to a safe filename."""
        safe = "".join(c if c.isalnum() or c in " -_" else "_" for c in title)
        safe = safe.strip().replace(" ", "_")[:60]
        return safe or "presentation"
