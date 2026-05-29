"""
Slide renderers — each slide type knows how to fill its layout placeholders.
"""

from __future__ import annotations

from typing import Optional, TYPE_CHECKING
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

if TYPE_CHECKING:
    from pptx import Presentation as PresentationType
    from pptx.slide import Slide
    from pptx.slide import SlideLayout


def _set_text(shape, text: str, font_name: Optional[str] = None,
              font_size: Optional[float] = None, color: Optional[str] = None,
              bold: Optional[bool] = None, alignment: Optional[str] = None):
    """Set text on a shape, preserving or overriding formatting."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = text

    if font_name:
        p.font.name = font_name
    if font_size:
        p.font.size = Pt(font_size)
    if color:
        try:
            p.font.color.rgb = RGBColor.from_string(color.lstrip("#"))
        except Exception:
            pass
    if bold is not None:
        p.font.bold = bold

    # Alignment
    align_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    if alignment and alignment.lower() in align_map:
        p.alignment = align_map[alignment.lower()]


def _set_bullets(tf, items: list[str], font_name: Optional[str] = None,
                 font_size: Optional[float] = None, color: Optional[str] = None):
    """Set bullet list on a text frame."""
    tf.clear()
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            from pptx.oxml.ns import qn
            p = tf.add_paragraph()

        p.text = item
        p.level = 0

        if font_name:
            p.font.name = font_name
        if font_size:
            p.font.size = Pt(font_size)
        if color:
            try:
                p.font.color.rgb = RGBColor.from_string(color.lstrip("#"))
            except Exception:
                pass


def render_title_slide(prs: PresentationType, layout: SlideLayout,
                       title: str, subtitle: Optional[str] = None,
                       theme_font: Optional[str] = None) -> Slide:
    """Render a title/cover slide."""
    slide = prs.slides.add_slide(layout)

    for ph in slide.placeholders:
        ph_type = ph.placeholder_format.type if ph.placeholder_format else ""
        try:
            if ph_type in ("title", "ctrTitle"):
                _set_text(ph, title, font_name=theme_font or ph.font.name,
                          font_size=44, bold=True)
            elif ph_type == "subTitle":
                _set_text(ph, subtitle or "", font_name=theme_font or ph.font.name,
                          font_size=20)
        except Exception:
            pass

    return slide


def render_section_slide(prs: PresentationType, layout: SlideLayout,
                         title: str, subtitle: Optional[str] = None,
                         theme_font: Optional[str] = None) -> Slide:
    """Render a section divider slide."""
    slide = prs.slides.add_slide(layout)

    for ph in slide.placeholders:
        ph_type = ph.placeholder_format.type if ph.placeholder_format else ""
        try:
            if ph_type in ("title", "ctrTitle"):
                _set_text(ph, title, font_name=theme_font or ph.font.name,
                          font_size=36, bold=True)
            elif ph_type == "subTitle":
                _set_text(ph, subtitle or "", font_name=theme_font or ph.font.name,
                          font_size=18)
            elif ph_type == "body":
                _set_text(ph, subtitle or "", font_name=theme_font or ph.font.name,
                          font_size=18)
        except Exception:
            pass

    return slide


def render_content_slide(prs: PresentationType, layout: SlideLayout,
                         title: str,
                         bullets: Optional[list[str]] = None,
                         body_text: Optional[str] = None,
                         theme_font: Optional[str] = None) -> Slide:
    """Render a standard content slide with title and body."""
    slide = prs.slides.add_slide(layout)

    for ph in slide.placeholders:
        ph_type = ph.placeholder_format.type if ph.placeholder_format else ""
        try:
            if ph_type in ("title", "ctrTitle"):
                _set_text(ph, title, font_name=theme_font or ph.font.name,
                          font_size=32, bold=True)
            elif ph_type in ("body", "text"):
                if bullets:
                    _set_bullets(ph.text_frame, bullets,
                                 font_name=theme_font or ph.font.name,
                                 font_size=18)
                elif body_text:
                    _set_text(ph, body_text, font_name=theme_font or ph.font.name,
                              font_size=18)
        except Exception:
            pass

    return slide


def render_quote_slide(prs: PresentationType, layout: SlideLayout,
                       quote_text: str, author: Optional[str] = None,
                       theme_font: Optional[str] = None) -> Slide:
    """Render a quote slide."""
    slide = prs.slides.add_slide(layout)

    # Try to find the quote placeholder or use body
    for ph in slide.placeholders:
        ph_type = ph.placeholder_format.type if ph.placeholder_format else ""
        try:
            if ph_type in ("title", "ctrTitle"):
                _set_text(ph, f"\"{quote_text}\"", font_name=theme_font or ph.font.name,
                          font_size=28)
                ph.text_frame.paragraphs[0].font.italic = True
            elif ph_type == "subTitle":
                _set_text(ph, f"— {author}" if author else "",
                          font_name=theme_font or ph.font.name,
                          font_size=16, color="#888888")
            elif ph_type in ("body", "text"):
                _set_text(ph, f"\"{quote_text}\"", font_name=theme_font or ph.font.name,
                          font_size=28)
                ph.text_frame.paragraphs[0].font.italic = True
        except Exception:
            pass

    return slide


def render_data_slide(prs: PresentationType, layout: SlideLayout,
                      title: str,
                      data_points: list[dict],
                      chart_type: str = "bar",
                      unit: Optional[str] = None,
                      theme_font: Optional[str] = None) -> Slide:
    """Render a data slide with numbers/chart info in placeholders."""
    slide = prs.slides.add_slide(layout)

    # Set title
    for ph in slide.placeholders:
        ph_type = ph.placeholder_format.type if ph.placeholder_format else ""
        try:
            if ph_type in ("title", "ctrTitle"):
                _set_text(ph, title, font_name=theme_font or ph.font.name,
                          font_size=32, bold=True)
        except Exception:
            pass

    # For data slides, format data points as text in body
    for ph in slide.placeholders:
        ph_type = ph.placeholder_format.type if ph.placeholder_format else ""
        try:
            if ph_type in ("body", "text", "object"):
                lines = []
                for pt in data_points:
                    label = pt.get("label", "")
                    value = pt.get("value", "")
                    if unit:
                        lines.append(f"▸ {label}: {value} {unit}")
                    else:
                        lines.append(f"▸ {label}: {value}")

                if chart_type in ("number", "big_number") and data_points:
                    # Show the first data point prominently
                    pt = data_points[0]
                    val = pt.get("value", "")
                    label = pt.get("label", "")
                    display = f"{val}{unit or ''}\n{label}"
                    _set_text(ph, display, font_name=theme_font or ph.font.name,
                              font_size=24)
                else:
                    _set_bullets(ph.text_frame, lines,
                                 font_name=theme_font or ph.font.name,
                                 font_size=18)
        except Exception:
            pass

    return slide


def render_end_slide(prs: PresentationType, layout: SlideLayout,
                     title: str = "Thank You",
                     subtitle: Optional[str] = None,
                     theme_font: Optional[str] = None) -> Slide:
    """Render a closing slide."""
    slide = prs.slides.add_slide(layout)

    for ph in slide.placeholders:
        ph_type = ph.placeholder_format.type if ph.placeholder_format else ""
        try:
            if ph_type in ("title", "ctrTitle"):
                _set_text(ph, title, font_name=theme_font or ph.font.name,
                          font_size=40, bold=True)
            elif ph_type == "subTitle":
                _set_text(ph, subtitle or "", font_name=theme_font or ph.font.name,
                          font_size=18)
        except Exception:
            pass

    return slide


def render_icons_slide(prs: PresentationType, layout: SlideLayout,
                       title: str = "Icons Reference",
                       dna=None,
                       theme_font: Optional[str] = None) -> Slide:
    """Render a slide listing all shapes/icons found in the template."""
    slide = prs.slides.add_slide(layout)

    # Title
    for ph in slide.placeholders:
        ph_type = ph.placeholder_format.type if ph.placeholder_format else ""
        try:
            if ph_type in ("title", "ctrTitle"):
                _set_text(ph, title, font_name=theme_font or ph.font.name,
                          font_size=32, bold=True)
            elif ph_type in ("body", "text") and dna:
                shapes_text = []
                for s in dna.template_shapes[:50]:
                    shapes_text.append(
                        f"  {s.name} ({s.shape_type}) — "
                        f"{s.width:.1f}×{s.height:.1f}in "
                        f"at ({s.left:.1f}, {s.top:.1f})"
                    )
                if shapes_text:
                    _set_text(ph, "\n".join(shapes_text),
                              font_name=theme_font or ph.font.name,
                              font_size=10)
        except Exception:
            pass

    return slide
