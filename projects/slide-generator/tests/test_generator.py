"""
Tests for the slide generator.
"""
import pytest
import yaml
from pathlib import Path

from src.plan_parser import parse_plan, parse_plan_string, Plan
from src.content_parser import parse_content, parse_content_string, Content


class TestPlanParser:
    """Tests for the plan parser."""

    def test_parse_basic_yaml(self, tmp_path):
        plan_yaml = """
title: "Test Presentation"
slides:
  - title: "Cover"
    type: title
  - title: "Intro"
    type: content
  - title: "Thanks"
    type: end
"""
        p = Path(tmp_path / "plan.yaml")
        p.write_text(plan_yaml)
        plan = parse_plan(p)
        assert plan.title == "Test Presentation"
        assert len(plan.slides) == 3
        assert plan.slides[0].type == "title"
        assert plan.slides[0].title == "Cover"
        assert plan.slides[2].type == "end"

    def test_parse_from_string(self):
        plan_yaml = """
title: "Test"
slides:
  - title: "Slide 1"
    type: content
"""
        plan = parse_plan_string(plan_yaml)
        # Auto-prepends a title slide
        assert len(plan.slides) == 2
        assert plan.slides[0].type == "title"
        assert plan.slides[1].type == "content"

    def test_auto_prepend_title(self):
        """If first slide isn't title, auto-prepend one."""
        plan_yaml = """
title: "My Pres"
slides:
  - title: "Intro"
    type: content
"""
        plan = parse_plan_string(plan_yaml)
        assert plan.slides[0].type == "title"
        assert plan.slides[0].title == "My Pres"
        assert plan.slides[1].type == "content"

    def test_invalid_type_raises(self):
        plan_yaml = """
slides:
  - title: "Bad"
    type: invalid_type
"""
        with pytest.raises(ValueError, match="invalid type"):
            parse_plan_string(plan_yaml)

    def test_empty_slides_raises(self):
        plan_yaml = """
title: "Empty"
slides: []
"""
        with pytest.raises(ValueError, match="at least one slide"):
            parse_plan_string(plan_yaml)

    def test_section_type(self):
        plan_yaml = """
title: "Test"
slides:
  - title: "Section 1"
    type: section
  - title: "Details"
    type: content
"""
        plan = parse_plan_string(plan_yaml)
        assert plan.slides[0].type == "title"  # auto-prepended
        assert plan.slides[1].type == "section"
        assert plan.slides[2].type == "content"


class TestContentParser:
    """Tests for the content parser."""

    def test_parse_content_text(self, tmp_path):
        content_yaml = """
slides:
  - slide_ref: "Intro"
    text:
      - style: body
        bullets:
          - "First point"
          - "Second point"
    notes: "Speaker note here"
"""
        p = Path(tmp_path / "content.yaml")
        p.write_text(content_yaml)
        content = parse_content(p)
        assert len(content.slides) == 1
        assert content.slides[0].slide_ref == "Intro"
        assert len(content.slides[0].text_blocks[0].bullets) == 2

    def test_parse_quotes(self):
        content_yaml = """
slides:
  - slide_ref: "Quote"
    quotes:
      - text: "To be or not to be"
        author: "Shakespeare"
"""
        content = parse_content_string(content_yaml)
        assert content.slides[0].quotes[0].text == "To be or not to be"
        assert content.slides[0].quotes[0].author == "Shakespeare"

    def test_parse_data(self):
        content_yaml = """
slides:
  - slide_ref: "Data"
    data:
      - title: "Revenue"
        chart_type: "bar"
        unit: "k€"
        points:
          - label: "Q1"
            value: 100
          - label: "Q2"
            value: 150
"""
        content = parse_content_string(content_yaml)
        assert content.slides[0].data[0].chart_type == "bar"
        assert len(content.slides[0].data[0].data_points) == 2

    def test_parse_lists(self):
        content_yaml = """
slides:
  - slide_ref: "Features"
    lists:
      - title: "Key Features"
        items:
          - "Feature A"
          - "Feature B"
"""
        content = parse_content_string(content_yaml)
        assert content.slides[0].lists[0].title == "Key Features"
        assert content.slides[0].lists[0].items == ["Feature A", "Feature B"]

    def test_multi_content_slide(self):
        """A slide can have text + quote + data."""
        content_yaml = """
slides:
  - slide_ref: "Mixed"
    text:
      - style: body
        text: "Intro text"
    quotes:
      - text: "Quote"
    data:
      - title: "Data"
        points:
          - label: "X"
            value: 1
"""
        content = parse_content_string(content_yaml)
        sc = content.slides[0]
        assert sc.text_blocks[0].text == "Intro text"
        assert sc.quotes[0].text == "Quote"
        assert sc.data[0].data_points[0]["label"] == "X"


class TestSlideBuilder:
    """Integration tests for the slide builder (requires a real .pptx template)."""

    def test_builder_requires_template(self):
        from src.slide_builder import SlideBuilder
        from pptx.exc import PackageNotFoundError
        with pytest.raises((FileNotFoundError, PackageNotFoundError)):
            SlideBuilder("/nonexistent/template.pptx")

    def test_build_with_minimal_plan(self, tmp_path):
        """Build with a real generated template and minimal plan."""
        from pptx import Presentation
        from pptx.util import Inches
        from src.slide_builder import SlideBuilder
        from src.plan_parser import parse_plan_string

        # Create a minimal template .pptx
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        prs.save(str(tmp_path / "minimal_template.pptx"))

        # Create a plan
        plan = parse_plan_string("""
title: "Test Build"
slides:
  - title: "Slide 1"
    type: content
""")

        builder = SlideBuilder(tmp_path / "minimal_template.pptx")
        output = builder.build(plan, output_path=tmp_path / "output.pptx")

        assert output.exists()
        assert output.suffix == ".pptx"

        # Verify output has slides
        result = Presentation(str(output))
        assert len(result.slides) >= 1


class TestTemplateParser:
    """Tests for the template parser."""

    def test_parse_minimal_template(self, tmp_path):
        from pptx import Presentation
        from pptx.util import Inches
        from src.template_parser import parse_template

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        prs.save(str(tmp_path / "test_template.pptx"))

        dna = parse_template(tmp_path / "test_template.pptx")
        assert dna.slide_width == pytest.approx(13.333, rel=0.01)
        assert dna.slide_height == pytest.approx(7.5, rel=0.01)

    def test_list_layouts(self, tmp_path):
        from pptx import Presentation
        from src.template_parser import list_layouts

        prs = Presentation()
        prs.save(str(tmp_path / "template.pptx"))

        layouts = list_layouts(tmp_path / "template.pptx")
        assert isinstance(layouts, list)
        # Should at least have the default blank layout
        assert len(layouts) > 0
